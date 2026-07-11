# -*- coding: utf-8 -*-
"""產生「台股分析工具 — 功能簡報」PowerPoint 檔（10 頁，深色交易終端風）。"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# ── palette ──
INK    = RGBColor(0x0d, 0x0e, 0x15)
PANEL  = RGBColor(0x17, 0x19, 0x22)
PANEL2 = RGBColor(0x1f, 0x22, 0x30)
LINE   = RGBColor(0x2a, 0x2d, 0x3d)
LINE2  = RGBColor(0x38, 0x3c, 0x50)
TEXT   = RGBColor(0xec, 0xee, 0xf5)
MUTED  = RGBColor(0x8b, 0x8f, 0xa6)
DIM    = RGBColor(0x56, 0x5a, 0x70)
GOLD   = RGBColor(0xf7, 0xb3, 0x4c)
UP     = RGBColor(0xff, 0x5d, 0x5d)   # 台股漲＝紅
DOWN   = RGBColor(0x35, 0xd6, 0xa0)   # 台股跌＝綠
BLUE   = RGBColor(0x6d, 0x8c, 0xff)
PURPLE = RGBColor(0xa0, 0x6c, 0xff)

SANS = "Microsoft JhengHei"
MONO = "Consolas"

EMU_IN = 914400
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = 13.333, 7.5
BLANK = prs.slide_layouts[6]

LM = 0.92          # 左右邊界
CW = SW - 2 * LM   # 內容寬


def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = INK
    r.line.fill.background()
    r.shadow.inherit = False
    return s


def txt(s, l, t, w, h, runs, size=16, color=TEXT, bold=False, font=SANS,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.0, letter=None):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    if isinstance(runs, str):
        runs = [(runs, {})]
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = spacing
    for text, opt in runs:
        r = p.add_run(); r.text = text
        f = r.font
        f.size = Pt(opt.get("size", size))
        f.name = opt.get("font", font)
        f.bold = opt.get("bold", bold)
        f.color.rgb = opt.get("color", color)
        if opt.get("letter", letter) is not None:
            _spc(r, opt.get("letter", letter))
    return tb


def _spc(run, pts):
    run.font._rPr.set("spc", str(int(pts * 100)))


def panel(s, l, t, w, h, fill=PANEL, line=LINE, radius=0.10, top_accent=None):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    if top_accent is not None:
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l + 0.14), Inches(t), Inches(w - 0.28), Inches(0.045))
        bar.fill.solid(); bar.fill.fore_color.rgb = top_accent
        bar.line.fill.background(); bar.shadow.inherit = False
    return sp


def eyebrow(s, num, tab):
    dash = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(LM), Inches(0.62), Inches(0.30), Inches(0.028))
    dash.fill.solid(); dash.fill.fore_color.rgb = GOLD
    dash.line.fill.background(); dash.shadow.inherit = False
    txt(s, LM + 0.42, 0.48, 8, 0.34,
        [(f"{num}", {"color": GOLD, "font": MONO, "size": 12.5, "letter": 2}),
         ("   " + tab, {"color": MUTED, "font": MONO, "size": 12.5, "letter": 2})],
        anchor=MSO_ANCHOR.MIDDLE)


def title(s, text, sub=None):
    txt(s, LM, 0.92, CW, 0.95, text, size=32, bold=True, color=TEXT)
    if sub:
        txt(s, LM, 1.92, CW, 0.75, sub, size=14.5, color=MUTED, spacing=1.15)


def footer(s, page):
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(LM), Inches(6.86), Inches(CW), Inches(0.012))
    ln.fill.solid(); ln.fill.fore_color.rgb = LINE
    ln.line.fill.background(); ln.shadow.inherit = False
    txt(s, LM, 6.98, 8, 0.3,
        [("台股分析工具", {"color": GOLD, "font": MONO, "size": 10.5, "bold": True, "letter": 1}),
         ("  · 功能簡報", {"color": DIM, "font": MONO, "size": 10.5, "letter": 1})],
        anchor=MSO_ANCHOR.MIDDLE)
    txt(s, SW - LM - 3, 6.98, 3, 0.3, f"{page:02d} / 10",
        font=MONO, size=10.5, color=DIM, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, letter=1)


# ══════════════ 1 · COVER ══════════════
import math
UP_DIM   = RGBColor(0x63, 0x30, 0x38)   # 暗紅（背景 K 線用，讓標題清楚）
DOWN_DIM = RGBColor(0x22, 0x4a, 0x40)   # 暗綠
s = slide()
# ambient candlesticks（暗色，當背景）
n = 46
price = SH * 0.55
cw = SW / n
for i in range(n):
    val = math.sin((i + .5) * 1.7) + math.sin((i + .5) * 0.6) * 1.4
    op = price
    price += val * (SH * 0.03)
    price = max(SH * 0.16, min(SH * 0.84, price))
    cl = price
    top = min(op, cl); bot = max(op, cl)
    col = UP_DIM if cl < op else DOWN_DIM
    cx = i * cw + cw / 2
    wick = 0.16 + abs(math.sin(i * 2.3)) * 0.18
    wk = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx - 0.01), Inches(top - wick), Inches(0.02), Inches((bot - top) + 2 * wick))
    wk.fill.solid(); wk.fill.fore_color.rgb = col; wk.line.fill.background(); wk.shadow.inherit = False
    bd = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx - cw * 0.28), Inches(top), Inches(cw * 0.56), Inches(max(0.05, bot - top)))
    bd.fill.solid(); bd.fill.fore_color.rgb = col; bd.line.fill.background(); bd.shadow.inherit = False

txt(s, 0, 2.05, SW, 0.5, "功能簡報 · 10 頁", font=MONO, size=13, color=GOLD,
    align=PP_ALIGN.CENTER, letter=6)
txt(s, 0, 2.55, SW, 1.6,
    [("台股", {"size": 74, "bold": True, "color": TEXT}),
     ("分析", {"size": 74, "bold": True, "color": UP}),
     ("工具", {"size": 74, "bold": True, "color": TEXT})],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(s, 0, 4.35, SW, 0.5, "看盤、警報、存股，一站整合的個人台股儀表板",
    size=18, color=MUTED, align=PP_ALIGN.CENTER)
# chips
chips = ["📈 圖表分析", "🔔 買賣點警鐘", "💰 存股", "iamlchang-max.github.io/stock-tw"]
cx = LM + 1.4
cy = 5.25
xs = 2.0
total_w = 0
widths = [1.9, 2.3, 1.5, 4.6]
gap = 0.25
start = (SW - (sum(widths) + gap * (len(widths) - 1))) / 2
x = start
for c, w in zip(chips, widths):
    ch = panel(s, x, cy, w, 0.52, fill=PANEL, line=LINE2, radius=0.5)
    txt(s, x, cy, w, 0.52, c, font=MONO, size=12.5,
        color=(GOLD if c.startswith("iaml") else TEXT),
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    x += w + gap

# ══════════════ 2 · ARCHITECTURE ══════════════
s = slide()
eyebrow(s, "01", "架構總覽")
title(s, "純前端 × Apps Script 後端",
      "前端部署在 GitHub Pages，資料與背景監控交給 Google Apps Script；行情走 Yahoo Finance，經自家代理避開 IP 封鎖。")
cy = 2.75; ch = 3.9
lw = CW * 0.56 - 0.15
rw = CW * 0.44 - 0.15
panel(s, LM, cy, lw, ch, fill=PANEL, line=LINE)
specs = [
    ("前端", "原生 HTML / CSS / JS", "index.html · app.css · app.js（GitHub Pages）"),
    ("後端", "Google Apps Script", "gas-code.gs — 資料儲存、行情代理、背景巡檢"),
    ("圖表", "Lightweight Charts", "K 線 · 均線 · 布林 · MACD · RSI · 量"),
    ("行情", "Yahoo Finance", "透過 GAS proxy（TWSE MIS 會擋 Google IP）"),
    ("儲存", "ScriptProperties + localStorage", "警報／持股存後端；設定與快取存本機"),
]
rh = (ch - 0.5) / len(specs)
for i, (k, v, sub) in enumerate(specs):
    ry = cy + 0.28 + i * rh
    txt(s, LM + 0.3, ry, 1.4, rh, k, font=MONO, size=12, color=GOLD, letter=1, anchor=MSO_ANCHOR.TOP)
    txt(s, LM + 1.7, ry - 0.02, lw - 2.0, rh,
        [(v, {"size": 14.5, "color": TEXT}), ("\n", {}), (sub, {"size": 11.5, "color": MUTED})],
        spacing=1.05)
    if i < len(specs) - 1:
        dl = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(LM + 0.3), Inches(ry + rh - 0.12), Inches(lw - 0.6), Inches(0.008))
        dl.fill.solid(); dl.fill.fore_color.rgb = LINE; dl.line.fill.background(); dl.shadow.inherit = False
# right data-flow panel
rx = LM + lw + 0.3
panel(s, rx, cy, rw, ch, fill=PANEL, line=LINE)
txt(s, rx + 0.28, cy + 0.22, rw - 0.5, 0.3, "DATA FLOW", font=MONO, size=11, color=MUTED, letter=2)
flow = [("瀏覽器（GitHub Pages）", "UI", TEXT),
        ("↕ fetch / POST", "JSON", UP),
        ("Apps Script Web App", "/exec", TEXT),
        ("↕ proxy", "Yahoo", TEXT),
        ("背景 cron（每分鐘）", "Telegram", GOLD)]
fy = cy + 0.72
for k, v, c in flow:
    txt(s, rx + 0.28, fy, rw - 1.7, 0.34, k, size=12.5, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, rx + rw - 1.6, fy, 1.35, 0.34, v, font=MONO, size=12.5, color=c, bold=True,
        align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    dl = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(rx + 0.28), Inches(fy + 0.36), Inches(rw - 0.56), Inches(0.008))
    dl.fill.solid(); dl.fill.fore_color.rgb = LINE; dl.line.fill.background(); dl.shadow.inherit = False
    fy += 0.5
tags = ["doGet", "doPost", "checkAlerts", "handleProxy"]
tx = rx + 0.28
for t in tags:
    w = 0.32 + len(t) * 0.095
    tp = panel(s, tx, fy + 0.05, w, 0.34, fill=PANEL2, line=LINE2, radius=0.3)
    txt(s, tx, fy + 0.05, w, 0.34, t, font=MONO, size=10.5, color=TEXT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tx += w + 0.14
footer(s, 2)

# ══════════════ 3 · THREE TABS ══════════════
s = slide()
eyebrow(s, "02", "分頁概覽")
title(s, "三個分頁，一條動線",
      "從看圖找標的、設好買賣點監控、再回頭盤點存股損益 —— 三頁彼此串接，代號到處可點跳圖。")
cards = [
    ("📈", "圖表分析", BLUE, "輸入代號或搜尋名稱，繪出完整技術線圖與即時資訊。",
     ["K 線＋均線＋布林", "MACD／RSI／量", "我的持股快選"]),
    ("🔔", "買賣點警鐘", GOLD, "預載持股，設定買點與賣點，觸價由後端發 Telegram 通知。",
     ["買賣雙價監控", "現價紅漲綠跌", "背景每分鐘巡檢"]),
    ("💰", "存股", UP, "盤點持股市值、損益與殖利率，圖表化投資組合分布。",
     ["總損益＋報酬率", "市場／個人殖利率", "標籤篩選與排序"]),
]
cw3 = (CW - 2 * 0.35) / 3
cy = 2.85; ch = 3.7
for i, (ic, h, ac, desc, pts) in enumerate(cards):
    x = LM + i * (cw3 + 0.35)
    panel(s, x, cy, cw3, ch, fill=PANEL, line=LINE, top_accent=ac)
    txt(s, x + 0.32, cy + 0.32, cw3 - 0.6, 0.6, ic, size=30)
    txt(s, x + 0.32, cy + 1.02, cw3 - 0.6, 0.5, h, size=21, bold=True, color=TEXT)
    txt(s, x + 0.32, cy + 1.62, cw3 - 0.64, 0.9, desc, size=13, color=MUTED, spacing=1.2)
    py = cy + 2.55
    for p in pts:
        txt(s, x + 0.32, py, cw3 - 0.6, 0.34,
            [("▸  ", {"color": ac, "size": 13}), (p, {"color": TEXT, "size": 13})],
            anchor=MSO_ANCHOR.MIDDLE)
        py += 0.36
footer(s, 3)

# ══════════════ 4 · CHART I ══════════════
s = slide()
eyebrow(s, "03", "圖表分析")
title(s, "查詢即繪圖，技術指標一次到位")
feats = [
    ("🔍", "代號 / 名稱查詢", "輸入代號或用中文名稱模糊搜尋，選 1～12 月區間，按查詢即繪圖。"),
    ("📊", "即時資訊面板", "名稱、股價、開高低收、昨收、成交量、漲跌，紅漲綠跌一眼判讀。"),
    ("📈", "四組技術圖表", "K 線＋移動平均＋布林通道，另附 MACD、RSI 與成交量。"),
]
cy = 2.5
lw = CW * 0.5 - 0.2
for i, (ic, h, p) in enumerate(feats):
    y = cy + i * 1.25
    ib = panel(s, LM, y, 0.55, 0.55, fill=PANEL2, line=LINE2, radius=0.28)
    txt(s, LM, y, 0.55, 0.55, ic, size=17, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, LM + 0.75, y - 0.04, lw - 0.75, 0.4, h, size=15.5, bold=True, color=TEXT)
    txt(s, LM + 0.75, y + 0.36, lw - 0.75, 0.7, p, size=12.5, color=MUTED, spacing=1.15)
# right mock
rx = LM + lw + 0.4
rw = CW - lw - 0.4
panel(s, rx, cy, rw, 4.0, fill=PANEL, line=LINE)
txt(s, rx + 0.3, cy + 0.24, rw - 2, 0.34, "2330 台積電", font=MONO, size=13, color=TEXT, bold=True)
txt(s, rx + rw - 2.3, cy + 0.24, 2.0, 0.34, "▲ 2415.00", font=MONO, size=13, color=UP, bold=True,
    align=PP_ALIGN.RIGHT)
rows = [("開盤", "2402.00", TEXT), ("最高 / 最低", "2420 / 2398", TEXT), ("漲跌", "+18.00 (+0.75%)", UP)]
ry = cy + 0.72
for k, v, c in rows:
    txt(s, rx + 0.3, ry, 2.5, 0.32, k, size=12.5, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, rx + rw - 3.0, ry, 2.7, 0.32, v, font=MONO, size=12.5, color=c, bold=True, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    ry += 0.42
# mini candles
mcx = rx + 0.3; mcw = rw - 0.6; mcy = cy + 2.15; mch = 1.15
mbg = panel(s, mcx, mcy, mcw, mch, fill=RGBColor(0x10, 0x12, 0x1b), line=LINE, radius=0.08)
nn = 26; price = mcy + mch * 0.5; ccw = mcw / nn
for i in range(nn):
    val = math.sin((i + 3.2) * 1.7) + math.sin((i + 3.2) * 0.6) * 1.3
    op = price; price += val * (mch * 0.12)
    price = max(mcy + 0.12, min(mcy + mch - 0.12, price)); cl = price
    top = min(op, cl); bot = max(op, cl); col = UP if cl < op else DOWN
    cxx = mcx + i * ccw + ccw / 2
    bd = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cxx - ccw * 0.28), Inches(top), Inches(ccw * 0.56), Inches(max(0.04, bot - top)))
    bd.fill.solid(); bd.fill.fore_color.rgb = col; bd.line.fill.background(); bd.shadow.inherit = False
# indicator tags
tags = ["K 線", "MA", "Bollinger", "MACD", "RSI", "Volume"]
tx = rx + 0.3; ty = cy + 3.45
for t in tags:
    w = 0.34 + len(t) * 0.11
    if tx + w > rx + rw - 0.3:
        break
    panel(s, tx, ty, w, 0.34, fill=PANEL2, line=LINE2, radius=0.25)
    txt(s, tx, ty, w, 0.34, t, font=MONO, size=10, color=TEXT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tx += w + 0.12
footer(s, 4)

# ══════════════ 5 · CHART II ══════════════
s = slide()
eyebrow(s, "04", "圖表分析")
title(s, "我的持股快選 · 一鍵看圖",
      "左側欄自動載入存股清單，點任一檔即在本頁繪圖；新增／刪除持股會即時同步，不必切分頁。")
cy = 2.75
lw = CW * 0.42 - 0.2
panel(s, LM, cy, lw, 3.9, fill=PANEL, line=LINE)
sb = panel(s, LM + 0.28, cy + 0.28, lw - 0.56, 0.44, fill=PANEL2, line=LINE2, radius=0.2)
txt(s, LM + 0.45, cy + 0.28, lw - 0.7, 0.44, "🔍  搜尋名稱或代號…", size=12, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)
txt(s, LM + 0.28, cy + 0.92, lw - 0.5, 0.3, "我的持股", font=MONO, size=11, color=MUTED, letter=2)
holds = [("2330", "台積電"), ("2892", "第一金"), ("4961", "天鈺"), ("2881", "富邦金"), ("6456", "GIS-KY")]
hy = cy + 1.32
for c, nm in holds:
    row = panel(s, LM + 0.28, hy, lw - 0.56, 0.42, fill=PANEL2, line=LINE2, radius=0.25)
    txt(s, LM + 0.48, hy, 1.2, 0.42, c, font=MONO, size=12.5, color=GOLD, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, LM + 1.55, hy, lw - 1.8, 0.42, nm, size=12.5, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)
    hy += 0.5
# right feature list
rx = LM + lw + 0.45
rw = CW - lw - 0.45
rfeats = [
    ("⚡", "點擊即繪圖", "持股晶片、存股頁「看圖」鈕、警鐘頁股名，三處都能一鍵跳到本頁繪圖。"),
    ("🔄", "自動同步", "存股頁增刪持股後，快選清單即時更新；同代號多帳戶只列一次。"),
    ("🗂️", "名稱對照快取", "台股代號↔中文對照表本機快取 24 小時，搜尋不必每次重抓。"),
]
for i, (ic, h, p) in enumerate(rfeats):
    y = cy + 0.15 + i * 1.2
    ib = panel(s, rx, y, 0.55, 0.55, fill=PANEL2, line=LINE2, radius=0.28)
    txt(s, rx, y, 0.55, 0.55, ic, size=17, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, rx + 0.75, y - 0.04, rw - 0.75, 0.4, h, size=15.5, bold=True, color=TEXT)
    txt(s, rx + 0.75, y + 0.36, rw - 0.75, 0.7, p, size=12.5, color=MUTED, spacing=1.15)
footer(s, 5)

# ══════════════ 6 · ALERT I (flow) ══════════════
s = slide()
eyebrow(s, "05", "買賣點警鐘")
title(s, "設好價格，剩下的交給背景",
      "前端只負責設定；真正的監控與通知在 Apps Script 上每分鐘自動執行，關掉網頁也照跑。")
steps = [
    ("STEP 1", "✏️", "設定買賣點", "對持股填買點（跌到）或賣點（漲到），按儲存。"),
    ("STEP 2", "⏱️", "背景巡檢", "GAS 每分鐘抓現價，僅交易時段 09:00–13:30。"),
    ("STEP 3", "📨", "Telegram 通知", "觸價立即推播訊息，含股名、現價與觸發方向。"),
]
cy = 3.0; ch = 3.0
sw = (CW - 2 * 0.9) / 3
for i, (sn, ic, h, p) in enumerate(steps):
    x = LM + i * (sw + 0.9)
    panel(s, x, cy, sw, ch, fill=PANEL, line=LINE)
    txt(s, x, cy + 0.32, sw, 0.3, sn, font=MONO, size=11.5, color=GOLD, align=PP_ALIGN.CENTER, letter=2)
    txt(s, x, cy + 0.72, sw, 0.7, ic, size=34, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x, cy + 1.62, sw, 0.4, h, size=17, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    txt(s, x + 0.25, cy + 2.12, sw - 0.5, 0.7, p, size=12.5, color=MUTED, align=PP_ALIGN.CENTER, spacing=1.15)
    if i < 2:
        ax = x + sw + 0.16
        txt(s, ax, cy, 0.58, ch, "→", size=28, color=LINE2, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
footer(s, 6)

# ── table helper ──
def make_table(s, l, t, w, headers, rows, col_w, aligns, page_h=0.44, row_h=0.44):
    nrows = len(rows) + 1; ncols = len(headers)
    gt = s.shapes.add_table(nrows, ncols, Inches(l), Inches(t), Inches(w), Inches(row_h * nrows)).table
    gt.first_row = False; gt.horz_banding = False
    # kill default style banding by setting fills manually
    tot = sum(col_w)
    for ci, cwd in enumerate(col_w):
        gt.columns[ci].width = Inches(w * cwd / tot)
    for ci, htext in enumerate(headers):
        c = gt.cell(0, ci); c.fill.solid(); c.fill.fore_color.rgb = PANEL
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.margin_left = Inches(0.1); c.margin_right = Inches(0.1); c.margin_top = 0; c.margin_bottom = 0
        p = c.text_frame.paragraphs[0]; p.alignment = aligns[ci]
        r = p.add_run(); r.text = htext
        r.font.size = Pt(11); r.font.name = MONO; r.font.color.rgb = MUTED; r.font.bold = False
    for ri, row in enumerate(rows):
        for ci, (val, col, mono, bold) in enumerate(row):
            c = gt.cell(ri + 1, ci); c.fill.solid()
            c.fill.fore_color.rgb = PANEL if ri % 2 == 0 else RGBColor(0x14, 0x16, 0x1f)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_left = Inches(0.1); c.margin_right = Inches(0.1); c.margin_top = 0; c.margin_bottom = 0
            p = c.text_frame.paragraphs[0]; p.alignment = aligns[ci]
            r = p.add_run(); r.text = val
            r.font.size = Pt(12); r.font.name = MONO if mono else SANS
            r.font.color.rgb = col; r.font.bold = bold
    return gt


# ══════════════ 7 · ALERT II (monitor table) ══════════════
s = slide()
eyebrow(s, "06", "買賣點警鐘")
title(s, "買賣雙價，一張表管完",
      "主畫面預設載入持股，一檔一列，買點與賣點可各自設定；改價再存會自動重新武裝。")
headers = ["代號", "名稱", "現價", "買點 ↓", "賣點 ↑", "狀態"]
aligns  = [PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.LEFT]
colw    = [1.1, 1.5, 1.3, 1.3, 1.3, 2.4]
rows = [
    [("2330", GOLD, True, True), ("台積電", TEXT, False, False), ("2415.00", UP, True, True),   ("1800", TEXT, True, False), ("—", DIM, True, False),  ("● 買點監控中", BLUE, False, False)],
    [("2892", GOLD, True, True), ("第一金", TEXT, False, False), ("33.60", DOWN, True, True),    ("28.0", TEXT, True, False), ("—", DIM, True, False),  ("● 買點觸發 27.7", GOLD, False, False)],
    [("2881", GOLD, True, True), ("富邦金", TEXT, False, False), ("124.50", UP, True, True),     ("90", TEXT, True, False),   ("150", TEXT, True, False),("● 買賣監控中", BLUE, False, False)],
    [("4961", GOLD, True, True), ("天鈺",   TEXT, False, False), ("176.00", DOWN, True, True),   ("—", DIM, True, False),     ("220", TEXT, True, False),("● 賣點監控中", BLUE, False, False)],
]
make_table(s, LM, 2.95, CW, headers, rows, colw, aligns, row_h=0.62)
footer(s, 7)

# ══════════════ 8 · SAVE I (cards + pie) ══════════════
s = slide()
eyebrow(s, "07", "存股")
title(s, "投資組合，一眼盤點")
cy = 2.45
cards = [("總市值", "2.6M", TEXT), ("總損益", "+218K (+9.1%)", UP), ("年股息估算", "84K", TEXT), ("平均殖利率", "3.2%", TEXT)]
cwd = (CW - 3 * 0.28) / 4
for i, (l, v, c) in enumerate(cards):
    x = LM + i * (cwd + 0.28)
    panel(s, x, cy, cwd, 1.35, fill=PANEL, line=LINE)
    txt(s, x + 0.24, cy + 0.24, cwd - 0.4, 0.3, l, size=12.5, color=MUTED, letter=1)
    txt(s, x + 0.24, cy + 0.62, cwd - 0.4, 0.55, v, font=MONO, size=(19 if len(v) > 8 else 24), bold=True, color=c, anchor=MSO_ANCHOR.MIDDLE)
# pie chart
py = cy + 1.65; ph = 2.55
panel(s, LM, py, CW, ph, fill=PANEL, line=LINE)
chart_data = CategoryChartData()
chart_data.categories = ["2892 第一金", "2330 台積電", "6456 GIS-KY", "4961 天鈺", "其他"]
chart_data.add_series("市值佔比", (57.0, 27.0, 7.9, 5.1, 3.0))
gframe = s.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(LM + 0.3), Inches(py + 0.25), Inches(3.4), Inches(ph - 0.5), chart_data)
chart = gframe.chart
chart.has_title = False
chart.has_legend = False
plot = chart.plots[0]
pie_cols = [BLUE, DOWN, GOLD, UP, PURPLE]
for i, pt in enumerate(plot.series[0].points):
    pt.format.fill.solid(); pt.format.fill.fore_color.rgb = pie_cols[i]
    pt.format.line.color.rgb = INK; pt.format.line.width = Pt(2)
# legend (manual)
leg = [("2892 第一金", "57.0%", BLUE), ("2330 台積電", "27.0%", DOWN), ("6456 GIS-KY", "7.9%", GOLD),
       ("4961 天鈺", "5.1%", UP), ("其他", "3.0%", PURPLE)]
lx = LM + 4.1; ly = py + 0.45
for nm, pct, c in leg:
    sw_ = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(lx), Inches(ly + 0.06), Inches(0.18), Inches(0.18))
    sw_.fill.solid(); sw_.fill.fore_color.rgb = c; sw_.line.fill.background(); sw_.shadow.inherit = False
    txt(s, lx + 0.32, ly, 3.2, 0.32, nm, size=13, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, lx + 3.4, ly, 1.3, 0.32, pct, font=MONO, size=13, color=MUTED, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    ly += 0.4
footer(s, 8)

# ══════════════ 9 · SAVE II (holdings table) ══════════════
s = slide()
eyebrow(s, "08", "存股")
title(s, "持股清單：損益與殖利率同框",
      "現價對比成本算出損益（金額＋%，賺紅虧綠），同時列出市場殖利率與個人殖利率；可標籤篩選、依欄排序、看總計列。")
headers = ["代號", "名稱", "成本", "現價", "市值", "損益", "市場殖利", "個人殖利"]
aligns  = [PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT]
colw    = [1.0, 1.3, 1.1, 1.1, 1.3, 2.2, 1.2, 1.2]
rows = [
    [("2881", GOLD, True, True), ("富邦金", TEXT, False, False), ("89.00", TEXT, True, False), ("124.50", TEXT, True, False), ("15,563", TEXT, True, False), ("+4,438 (+39.9%)", UP, True, True), ("3.41%", TEXT, True, False), ("4.78%", TEXT, True, False)],
    [("2317", GOLD, True, True), ("鴻海", TEXT, False, False),   ("294.00", TEXT, True, False), ("237.50", TEXT, True, False), ("23,750", TEXT, True, False), ("−5,650 (−19.2%)", DOWN, True, True), ("3.02%", TEXT, True, False), ("2.44%", TEXT, True, False)],
    [("2330", GOLD, True, True), ("台積電", TEXT, False, False), ("610.00", TEXT, True, False), ("1050.0", TEXT, True, False), ("1,050K", TEXT, True, False), ("+440K (+72%)", UP, True, True), ("1.10%", TEXT, True, False), ("1.90%", TEXT, True, False)],
    [("2892", GOLD, True, True), ("第一金", TEXT, False, False), ("28.19", TEXT, True, False), ("33.60", TEXT, True, False), ("508,704", TEXT, True, False), ("+81,900 (+19%)", UP, True, True), ("2.83%", TEXT, True, False), ("3.37%", TEXT, True, False)],
]
make_table(s, LM, 3.05, CW, headers, rows, colw, aligns, row_h=0.6)
footer(s, 9)

# ══════════════ 10 · BACKEND & CLOSE ══════════════
s = slide()
eyebrow(s, "09", "後端 · 部署")
title(s, "背景引擎與上線方式")
cy = 2.6
lw = CW * 0.54 - 0.2
back = [
    ("資料 API", "doGet 回傳 alerts＋holdings；doPost 支援警報／持股 CRUD"),
    ("行情代理", "?proxy=url，僅白名單網域（TWSE／Yahoo）"),
    ("背景巡檢", "checkAlerts 每分鐘跑，僅交易時段觸價發通知"),
    ("前端部署", "git push → GitHub Pages 約 1–2 分鐘更新"),
    ("後端部署", "改 gas-code.gs 需重新部署；一般 CRUD 免動"),
]
rh = 0.82
for i, (k, v) in enumerate(back):
    y = cy + i * rh
    txt(s, LM, y + 0.04, 1.6, 0.5, k, font=MONO, size=12.5, color=GOLD, letter=1)
    txt(s, LM + 1.7, y, lw - 1.7, 0.7, v, size=13.5, color=TEXT, spacing=1.1)
    if i < len(back) - 1:
        dl = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(LM), Inches(y + rh - 0.16), Inches(lw), Inches(0.008))
        dl.fill.solid(); dl.fill.fore_color.rgb = LINE; dl.line.fill.background(); dl.shadow.inherit = False
# right closing card
rx = LM + lw + 0.4
rw = CW - lw - 0.4
panel(s, rx, cy, rw, 4.15, fill=PANEL, line=LINE)
txt(s, rx + 0.4, cy + 0.55, rw - 0.8, 1.4,
    [("看盤 · 警報 · 存股", {"size": 26, "bold": True, "color": TEXT}), ("\n", {}),
     ("都在一頁。", {"size": 26, "bold": True, "color": TEXT})], spacing=1.15)
txt(s, rx + 0.4, cy + 2.1, rw - 0.8, 0.9,
    [("純前端 + Apps Script · 無伺服器成本", {"size": 13, "color": MUTED}), ("\n", {}),
     ("資料自持、Telegram 即時通知", {"size": 13, "color": MUTED})], spacing=1.3)
txt(s, rx + 0.4, cy + 3.35, rw - 0.8, 0.5, "iamlchang-max.github.io/stock-tw",
    font=MONO, size=13.5, color=GOLD)
footer(s, 10)

prs.save("台股分析工具_功能簡報.pptx")
print("OK saved 台股分析工具_功能簡報.pptx  slides:", len(prs.slides._sldIdLst))
